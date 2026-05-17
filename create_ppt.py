from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette (exact web tokens) ──────────────────────────
PRIMARY   = RGBColor(0xBA, 0xD3, 0x31)   # #bad331  lime
DARK      = RGBColor(0x14, 0x3A, 0x3C)   # #143a3c  brand-dark teal
BLACK     = RGBColor(0x02, 0x06, 0x17)   # #020617
SLATE     = RGBColor(0x64, 0x74, 0x8B)   # #64748b
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF6, 0xF8)
BORDER    = RGBColor(0xF1, 0xF5, 0xF9)
DARK2     = RGBColor(0x0F, 0x2B, 0x2D)   # deeper teal bg card
TEAL_MID  = RGBColor(0x1E, 0x50, 0x52)
TEAL_TEXT = RGBColor(0x9B, 0xC2, 0xC4)
GREEN_OK  = RGBColor(0x14, 0x7A, 0x3C)

F_HEAD = "Helvetica Neue"
F_BODY = "Helvetica Neue"
BLANK  = 6

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ── Primitives ────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=13, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, font=F_BODY, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

def section_label(slide, text, l=0.4, t=0.28, color=SLATE):
    txt(slide, text, l, t, 8, 0.3, size=8.5, bold=True,
        color=color, font=F_HEAD)

def heading(slide, text, l, t, w, h, size=32, color=DARK, center=False):
    txt(slide, text, l, t, w, h, size=size, bold=True, color=color,
        font=F_HEAD, align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT)

def foot(slide):
    rect(slide, 0, 7.22, 13.33, 0.28, DARK)
    txt(slide, "INGATE  —  ingate.id", 0.5, 7.24, 6, 0.22,
        size=8, color=TEAL_TEXT, font=F_HEAD)


# ══════════════════════════════════════════════════════════════
# SLIDE 1  —  Cover
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])

rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)   # lime left stripe

# Decorative circles
c = s.shapes.add_shape(9, Inches(8.0), Inches(-1.5), Inches(7), Inches(7))
c.fill.solid(); c.fill.fore_color.rgb = DARK2; c.line.fill.background()
c2 = s.shapes.add_shape(9, Inches(10), Inches(4.2), Inches(4.5), Inches(4.5))
c2.fill.solid(); c2.fill.fore_color.rgb = DARK2; c2.line.fill.background()

txt(s, "INGATE", 0.5, 0.42, 5, 0.75,
    size=24, bold=True, color=PRIMARY, font=F_HEAD)
txt(s, "TICKETING PLATFORM  ·  INDONESIA",
    0.5, 1.28, 7, 0.32, size=8.5, bold=True,
    color=TEAL_TEXT, font=F_HEAD)

heading(s, "Jual Tiket Event Anda.", 0.5, 2.0, 11, 0.95, size=56, color=WHITE)
heading(s, "Online. Mudah.", 0.5, 2.95, 11, 0.9, size=56, color=WHITE)
heading(s, "Sekarang.", 0.5, 3.9, 11, 0.9, size=56, color=PRIMARY)

txt(s, "Platform ticketing online untuk event organizer Indonesia —\nbayar hanya saat tiket terjual, tanpa biaya di muka.",
    0.5, 5.05, 9, 0.9, size=14, color=TEAL_TEXT)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 2  —  Problem  (2×2 grid)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "01  /  MASALAH")
heading(s, "Masih Jual Tiket Cara Lama?", 0.4, 0.58, 12, 0.72, size=34, color=DARK)

problems = [
    ("01", "Transfer Manual & Konfirmasi WA",
     "Pembeli kirim bukti transfer, tim rekap satu per satu lewat WhatsApp — "
     "rawan salah catat, lambat dikonfirmasi, dan sama sekali tidak profesional "
     "di mata pembeli."),
    ("02", "Penukaran Tiket Kacau di Pintu Masuk",
     "Tanpa sistem digital, panitia kesulitan menukar e-ticket ke tiket gelang "
     "dengan cepat dan akurat. Antrian panjang, tiket palsu lolos, "
     "dan pengalaman tamu di menit pertama langsung rusak."),
    ("03", "Buta Data Penjualan Real-time",
     "Tidak tahu berapa tiket terjual hari ini, tipe mana yang laris, "
     "dan total pendapatan sudah berapa — semua harus rekap spreadsheet dulu "
     "sebelum bisa ambil keputusan."),
    ("04", "Kirim Tiket Satu Per Satu ke Pembeli",
     "Setelah pembayaran dikonfirmasi, tim harus kirim tiket manual ke email "
     "pembeli satu per satu — memakan waktu berjam-jam dan rawan terlewat "
     "saat event besar dengan ratusan pembeli."),
]

# 2×2 grid — card size 6.1 × 2.5
card_positions = [(0.38, 1.52), (6.85, 1.52), (0.38, 4.22), (6.85, 4.22)]
CW, CH = 6.1, 2.45

for i, ((lx, ly), (num, title, desc)) in enumerate(zip(card_positions, problems)):
    rect(s, lx, ly, CW, CH, WHITE)
    # left number strip
    strip_col = PRIMARY if i % 2 == 0 else DARK
    rect(s, lx, ly, 0.62, CH, strip_col)
    txt(s, num, lx + 0.05, ly + 0.82, 0.52, 0.7,
        size=20, bold=True,
        color=DARK if i % 2 == 0 else WHITE,
        font=F_HEAD, align=PP_ALIGN.CENTER)
    # title
    txt(s, title, lx + 0.82, ly + 0.18, 5.1, 0.48,
        size=13.5, bold=True, color=DARK, font=F_HEAD)
    # thin rule
    rect(s, lx + 0.82, ly + 0.72, 5.1, 0.03, BORDER)
    # description
    txt(s, desc, lx + 0.82, ly + 0.85, 5.1, 1.5,
        size=11.5, color=SLATE)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 3  —  Solusi
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "02  /  SOLUSI", color=TEAL_TEXT)
heading(s, "Satu Platform.\nSemua Beres.", 0.4, 0.62, 9, 1.2, size=40, color=WHITE)

txt(s, "Organizer cukup daftar, buat event, dan share link.\n"
       "INGATE yang urus pembayaran, e-ticket, monitoring penjualan, hingga penukaran tiket gelang di pintu masuk.",
    0.4, 1.95, 11.5, 0.85, size=14, color=TEAL_TEXT)

# Flow bar
flow_items = [
    "Daftar &\nBuat Event",
    "Atur Tiket\n& Kuota",
    "Share\nLink Event",
    "Pembeli Beli\n& Bayar",
    "E-Ticket\nOtomatis",
    "Tukar ke\nTiket Gelang",
]

BW = 1.75   # block width
AW = 0.30   # arrow block width
GAP = 0.04  # gap each side of arrow
STEP = BW + GAP + AW + GAP   # 2.13 per step

for i, label in enumerate(flow_items):
    lx = 0.38 + i * STEP
    bg = PRIMARY if i in (0, 2, 4) else TEAL_MID
    rect(s, lx, 3.0, BW, 1.55, bg)
    txt(s, label, lx + 0.08, 3.2, BW - 0.12, 1.1,
        size=12.5, bold=True,
        color=DARK if i in (0, 2, 4) else WHITE,
        font=F_HEAD, align=PP_ALIGN.CENTER)
    # Solid arrow block between steps
    if i < 5:
        ax = lx + BW + GAP
        rect(s, ax, 3.22, AW, 1.1, DARK)
        txt(s, "→", ax, 3.26, AW, 0.98,
            size=18, bold=True, color=PRIMARY,
            font=F_HEAD, align=PP_ALIGN.CENTER)

# Value props below flow
props = [
    ("Tanpa biaya setup", "Mulai gratis, bayar hanya saat tiket terjual."),
    ("Pembayaran otomatis", "QRIS, bank transfer, kartu kredit — semua terintegrasi."),
    ("E-ticket instan", "Email konfirmasi + QR Code dikirim otomatis ke pembeli."),
]
for i, (title, desc) in enumerate(props):
    lx = 0.4 + i * 4.3
    rect(s, lx, 4.85, 4.0, 1.9, DARK2)
    rect(s, lx, 4.85, 4.0, 0.06, PRIMARY)
    txt(s, title, lx + 0.2, 5.0, 3.6, 0.45,
        size=13, bold=True, color=PRIMARY, font=F_HEAD)
    txt(s, desc, lx + 0.2, 5.52, 3.6, 1.1, size=11.5, color=TEAL_TEXT)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 4  —  Cara Kerja  (timeline horizontal strip)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, DARK)

# ── Header kompak ────────────────────────────────────────────
section_label(s, "03  /  CARA KERJA", color=TEAL_MID)
txt(s, "Mulai dalam 4 Langkah.", 0.4, 0.55, 10, 0.72,
    size=32, bold=True, color=WHITE, font=F_HEAD)

# ── Layout constants ─────────────────────────────────────────
NUM_CX   = 0.95      # center-x kotak nomor
NUM_SZ   = 0.70      # ukuran kotak nomor
BH       = 1.35      # tinggi setiap band
BY_START = 1.48      # y mulai band pertama
GAP      = 0.06      # jarak antar band

# Spine vertikal — dari bawah kotak pertama ke atas kotak terakhir
SQ_TOP0  = BY_START + (BH - NUM_SZ) / 2
SQ_BOT0  = SQ_TOP0 + NUM_SZ
SQ_TOP3  = BY_START + 3 * BH + 3 * GAP + (BH - NUM_SZ) / 2
rect(s, NUM_CX - 0.03, SQ_BOT0, 0.06, SQ_TOP3 - SQ_BOT0, TEAL_MID)

steps = [
    ("01", "Daftar & Buat Event",
     "Isi nama event, tanggal, lokasi, dan deskripsi. "
     "Tambahkan banner untuk tampilan yang profesional."),
    ("02", "Tambah Tipe Tiket",
     "Buat tiket VIP, Regular, Early Bird — atur harga, kuota, "
     "dan jadwal buka-tutup penjualan per tipe."),
    ("03", "Share Link Event",
     "INGATE otomatis membuat halaman event publik. "
     "Share link ke media sosial, poster, atau grup komunitas."),
    ("04", "Pantau Penjualan & Tukar Tiket Gelang",
     "Monitor tiket terjual & pendapatan real-time. "
     "Di hari H, tim tukar e-ticket QR ke tiket gelang fisik."),
]

# Lebar kolom — tidak ada watermark agar tidak tumpang tindih
TITLE_X = 1.9;   TITLE_W = 4.3    # judul  x=1.9 .. 6.2
SEP_X   = 6.35                     # garis vertikal pemisah
DESC_X  = 6.6;   DESC_W  = 6.35   # deskripsi x=6.6 .. 12.95

for i, (num, title, desc) in enumerate(steps):
    by = BY_START + i * (BH + GAP)

    # Band background bergantian
    band_bg = RGBColor(0x10, 0x30, 0x32) if i % 2 == 0 else DARK
    rect(s, 0, by, 13.33, BH, band_bg)

    # Kotak nomor
    sq_top = by + (BH - NUM_SZ) / 2
    sq_col = PRIMARY if i % 2 == 0 else TEAL_MID
    rect(s, NUM_CX - NUM_SZ / 2, sq_top, NUM_SZ, NUM_SZ, sq_col)
    txt(s, num,
        NUM_CX - NUM_SZ / 2, sq_top,
        NUM_SZ, NUM_SZ,
        size=17, bold=True,
        color=DARK if i % 2 == 0 else WHITE,
        font=F_HEAD, align=PP_ALIGN.CENTER)

    # Konektor horizontal pendek (nomor → judul)
    rect(s, NUM_CX + NUM_SZ / 2 + 0.05, by + BH / 2 - 0.02,
         0.35, 0.04, sq_col)

    # Judul — center vertikal di dalam band
    txt(s, title, TITLE_X, by + (BH - 0.52) / 2,
        TITLE_W, 0.52,
        size=14, bold=True, color=WHITE, font=F_HEAD)

    # Separator vertikal
    rect(s, SEP_X, by + 0.18, 0.04, BH - 0.36,
         RGBColor(0x1E, 0x50, 0x52))

    # Deskripsi — hanya di kanan
    txt(s, desc, DESC_X, by + (BH - 0.65) / 2,
        DESC_W, 0.65,
        size=12, color=TEAL_TEXT)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 5  —  Fitur untuk Organizer
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "04  /  FITUR", color=TEAL_TEXT)
heading(s, "Semua yang Anda Butuhkan\nAda di Satu Tempat.", 0.4, 0.62, 10, 1.1, size=32, color=WHITE)

features = [
    ("Kelola Event",
     "Buat, edit, dan nonaktifkan event kapan saja. "
     "Atur status event: draft, aktif, atau selesai."),
    ("Multi-Tipe Tiket",
     "VIP, Regular, Early Bird, dan lainnya — "
     "tiap tipe punya harga, kuota, dan jadwal sendiri."),
    ("Halaman Event Publik",
     "Setiap event punya halaman profesional yang "
     "bisa langsung di-share dan diakses pembeli."),
    ("Pembayaran Terintegrasi",
     "QRIS, transfer bank, dan kartu kredit — "
     "semua konfirmasi otomatis, tanpa rekap manual."),
    ("E-Ticket Otomatis",
     "Setelah bayar, pembeli langsung terima "
     "e-ticket ber-QR Code di email mereka."),
    ("Dashboard Real-time",
     "Pantau total pendapatan, tiket terjual, "
     "dan event aktif — semua dalam satu layar."),
    ("Laporan Transaksi",
     "Riwayat lengkap semua transaksi: nama pembeli, "
     "tipe tiket, status, dan nominal — bisa difilter."),
    ("Profil Organizer",
     "Upload logo, isi bio, dan kelola informasi "
     "brand organizer Anda di platform INGATE."),
]

for i, (title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    lx = 0.4 + col * 3.25
    ly = 2.0 + row * 2.45

    rect(s, lx, ly, 3.0, 2.2, DARK2)
    rect(s, lx, ly, 3.0, 0.05, PRIMARY if row == 0 else TEAL_MID)

    txt(s, title, lx + 0.18, ly + 0.18, 2.65, 0.45,
        size=12.5, bold=True, color=PRIMARY, font=F_HEAD)
    txt(s, desc, lx + 0.18, ly + 0.72, 2.65, 1.3,
        size=11, color=TEAL_TEXT)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 6  —  Pengalaman Pembeli
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "05  /  PENGALAMAN PEMBELI")
heading(s, "Pembeli Anda Mendapat\nPengalaman Terbaik.", 0.4, 0.62, 9, 1.1, size=34, color=DARK)

txt(s, "Kesan pertama pembeli terhadap event Anda dimulai dari cara mereka membeli tiket.",
    0.4, 1.82, 10, 0.45, size=13, color=SLATE)

buyer_steps = [
    ("Halaman Event\nProfesional",
     "Pembeli melihat halaman event yang bersih, "
     "informatif, dan mobile-friendly."),
    ("Checkout\nMudah",
     "Isi data diri, pilih tipe & jumlah tiket, "
     "lanjut ke pembayaran — semua dalam hitungan menit."),
    ("Bayar via\nMidtrans",
     "QRIS, virtual account, kartu kredit. "
     "Konfirmasi instan, tidak perlu kirim bukti transfer."),
    ("E-Ticket\nLangsung di Email",
     "QR Code unik per tiket dikirim otomatis. "
     "Bisa disimpan di HP, tidak perlu cetak."),
    ("Tukar E-Ticket\nke Tiket Gelang",
     "Panitia scan QR Code pembeli, sistem verifikasi otomatis, "
     "lalu gelang fisik langsung diberikan. Cepat, anti-fraud."),
]

for i, (title, desc) in enumerate(buyer_steps):
    lx = 0.4 + i * 2.6
    ly = 2.55

    rect(s, lx, ly, 2.4, 4.35, WHITE)
    rect(s, lx, ly, 2.4, 0.06, PRIMARY if i % 2 == 0 else DARK)

    txt(s, str(i + 1), lx + 0.18, ly + 0.2, 0.45, 0.5,
        size=22, bold=True, color=BORDER, font=F_HEAD)
    txt(s, title, lx + 0.18, ly + 0.82, 2.05, 0.68,
        size=12.5, bold=True, color=DARK, font=F_HEAD)
    txt(s, desc, lx + 0.18, ly + 1.65, 2.05, 2.5,
        size=11.5, color=SLATE)

    if i < 4:
        txt(s, "→", lx + 2.42, ly + 1.95, 0.22, 0.45,
            size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 7  —  Monitoring Penjualan
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "06  /  MONITORING PENJUALAN", color=TEAL_TEXT)
heading(s, "Pantau Semua Penjualan\nSecara Real-time.", 0.4, 0.62, 10, 1.05, size=32, color=WHITE)

# ── Top stat cards (3 KPI) ────────────────────────────────────
stats = [
    ("Total Pendapatan",  "Rp 48.750.000",  "+12% dari event sebelumnya"),
    ("Tiket Terjual",     "1.250 / 2.000",  "62.5% kuota terisi"),
    ("Tiket Ditukar",     "380 gelang",     "Sudah ditukar di pintu masuk"),
]
for i, (label, val, note) in enumerate(stats):
    lx = 0.38 + i * 4.32
    rect(s, lx, 1.9, 4.05, 1.6, DARK2)
    rect(s, lx, 1.9, 4.05, 0.06, PRIMARY if i == 0 else TEAL_MID)
    txt(s, label.upper(), lx + 0.2, 2.04, 3.65, 0.32,
        size=8.5, bold=True, color=TEAL_TEXT, font=F_HEAD)
    txt(s, val, lx + 0.2, 2.4, 3.65, 0.62,
        size=21, bold=True, color=WHITE, font=F_HEAD)
    txt(s, note, lx + 0.2, 3.08, 3.65, 0.32,
        size=9.5, color=PRIMARY if i == 0 else TEAL_MID)

# ── Left: Riwayat Transaksi mockup ────────────────────────────
rect(s, 0.38, 3.7, 6.6, 3.28, DARK2)
rect(s, 0.38, 3.7, 6.6, 0.06, PRIMARY)
txt(s, "RIWAYAT TRANSAKSI", 0.55, 3.82, 4, 0.32,
    size=8.5, bold=True, color=TEAL_TEXT, font=F_HEAD)

# Column headers
th_cols = ["Pembeli", "Tiket", "Status", "Nominal"]
th_x    = [0.55, 2.3, 4.0, 5.1]
th_w    = [1.65, 1.6, 1.0, 1.6]
for lbl, cx, cw in zip(th_cols, th_x, th_w):
    txt(s, lbl, cx, 4.18, cw, 0.3,
        size=9, bold=True, color=TEAL_MID, font=F_HEAD)

rect(s, 0.55, 4.5, 6.25, 0.03, TEAL_MID)

# Mockup rows
mock_rows = [
    ("Budi Santoso",   "VIP · 2 tiket",  "PAID",    "Rp 600.000"),
    ("Sari Dewi",      "Regular · 1",    "PAID",    "Rp 150.000"),
    ("Reza Pratama",   "Early Bird · 3", "PENDING", "Rp 360.000"),
    ("Lina Kusuma",    "VIP · 1 tiket",  "PAID",    "Rp 300.000"),
]
STATUS_COLOR = {"PAID": GREEN_OK, "PENDING": RGBColor(0xE6, 0x8A, 0x1A)}
for r, (name, ticket, status, amount) in enumerate(mock_rows):
    ry = 4.56 + r * 0.58
    row_bg = DARK if r % 2 == 0 else DARK2
    rect(s, 0.38, ry, 6.6, 0.56, row_bg)
    vals = [name, ticket, status, amount]
    colors = [WHITE, TEAL_TEXT, STATUS_COLOR.get(status, WHITE), WHITE]
    bolds  = [False, False, True, False]
    for val, cx, cw, col, bd in zip(vals, th_x, th_w, colors, bolds):
        txt(s, val, cx, ry + 0.1, cw, 0.38,
            size=10.5, bold=bd, color=col, font=F_BODY)

# ── Right: Breakdown per Tipe Tiket ───────────────────────────
rect(s, 7.22, 3.7, 5.75, 3.28, DARK2)
rect(s, 7.22, 3.7, 5.75, 0.06, TEAL_MID)
txt(s, "PENJUALAN PER TIPE TIKET", 7.4, 3.82, 5.4, 0.32,
    size=8.5, bold=True, color=TEAL_TEXT, font=F_HEAD)

# Ticket type rows with progress bar
ticket_types = [
    ("VIP",        450,  500,  "Rp 135.000.000"),
    ("Regular",    680,  1000, "Rp 102.000.000"),
    ("Early Bird", 120,  500,  "Rp 12.000.000"),
]
BAR_W = 3.8
for i, (ttype, sold, total, revenue) in enumerate(ticket_types):
    ty = 4.3 + i * 0.95
    pct = sold / total

    txt(s, ttype, 7.4, ty, 2.0, 0.32,
        size=11, bold=True, color=WHITE, font=F_HEAD)
    txt(s, f"{sold}/{total} terjual", 9.45, ty, 2.0, 0.32,
        size=10, color=TEAL_TEXT, align=PP_ALIGN.RIGHT)

    # Progress bar background
    rect(s, 7.4, ty + 0.38, BAR_W, 0.22, RGBColor(0x1E, 0x40, 0x42))
    # Progress bar fill
    fill_w = max(BAR_W * pct, 0.15)
    bar_col = PRIMARY if pct > 0.5 else RGBColor(0xE6, 0x8A, 0x1A)
    rect(s, 7.4, ty + 0.38, fill_w, 0.22, bar_col)
    # Percentage label
    txt(s, f"{int(pct*100)}%", 11.22, ty + 0.36, 0.6, 0.26,
        size=9.5, bold=True, color=bar_col, font=F_HEAD)

    txt(s, revenue, 7.4, ty + 0.65, BAR_W + 0.8, 0.28,
        size=9.5, color=TEAL_MID)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 8  —  Pricing
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, OFF_WHITE)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "07  /  HARGA")
heading(s, "Transparan.\nBayar Saat Tiket Terjual.", 0.4, 0.58, 9, 1.05, size=34, color=DARK)

# ── Left column: 5% model card ────────────────────────────────
rect(s, 0.38, 1.82, 5.9, 5.1, DARK)
rect(s, 0.38, 1.82, 5.9, 0.07, PRIMARY)

txt(s, "MODEL KOMISI", 0.58, 1.97, 5.5, 0.35,
    size=8.5, bold=True, color=TEAL_TEXT, font=F_HEAD)

txt(s, "5%", 0.48, 2.35, 5.7, 1.35,
    size=88, bold=True, color=PRIMARY, font=F_HEAD, align=PP_ALIGN.CENTER)

txt(s, "dari harga tiket yang terjual", 0.58, 3.7, 5.5, 0.42,
    size=14, color=WHITE, font=F_BODY, align=PP_ALIGN.CENTER)

rect(s, 0.75, 4.25, 5.15, 0.03, TEAL_MID)

key_points = [
    ("Tanpa biaya di muka",  "Kami hanya terbayar saat tiket Anda terjual."),
    ("Admin Fee",            "Dapat didiskusikan bersama tim kami."),
]
for i, (title, desc) in enumerate(key_points):
    ty = 4.4 + i * 0.82
    txt(s, title, 0.58, ty, 5.5, 0.36,
        size=11.5, bold=True, color=PRIMARY, font=F_HEAD, align=PP_ALIGN.CENTER)
    txt(s, desc, 0.58, ty + 0.38, 5.5, 0.35,
        size=10.5, color=TEAL_TEXT, align=PP_ALIGN.CENTER)

# ── Right column: simulation table ────────────────────────────
rect(s, 6.75, 1.82, 6.2, 5.1, WHITE)
rect(s, 6.75, 1.82, 6.2, 0.07, DARK)

txt(s, "SIMULASI SEDERHANA", 6.95, 1.97, 5.8, 0.35,
    size=8.5, bold=True, color=DARK, font=F_HEAD)

# Table header
rect(s, 6.75, 2.42, 6.2, 0.48, DARK)
col_labels = ["Harga Tiket", "Platform Fee (5%)", "Anda Terima"]
col_x = [6.88, 8.78, 10.68]
col_w = [1.75, 1.9, 1.8]
for lbl, cx, cw in zip(col_labels, col_x, col_w):
    txt(s, lbl, cx, 2.46, cw, 0.38,
        size=9.5, bold=True, color=WHITE, font=F_HEAD, align=PP_ALIGN.CENTER)

# Simulation rows
sim_rows = [
    ("Rp 50.000",   "Rp 2.500",   "Rp 47.500"),
    ("Rp 100.000",  "Rp 5.000",   "Rp 95.000"),
    ("Rp 200.000",  "Rp 10.000",  "Rp 190.000"),
    ("Rp 500.000",  "Rp 25.000",  "Rp 475.000"),
]
RH = 0.62
for r, (price, fee, net) in enumerate(sim_rows):
    ry = 2.92 + r * RH
    row_bg = OFF_WHITE if r % 2 == 0 else WHITE
    rect(s, 6.75, ry, 6.2, RH - 0.02, row_bg)

    vals = [price, fee, net]
    colors = [BLACK, RGBColor(0xC0, 0x30, 0x20), GREEN_OK]
    for val, cx, cw, col in zip(vals, col_x, col_w, colors):
        txt(s, val, cx, ry + 0.1, cw, RH - 0.15,
            size=12, bold=(col == GREEN_OK),
            color=col, font=F_HEAD, align=PP_ALIGN.CENTER)

# Separator + note
rect(s, 6.75, 2.92 + 4 * RH + 0.05, 6.2, 0.03, BORDER)
txt(s, "* Belum termasuk Admin Fee (dapat didiskusikan)",
    6.85, 2.92 + 4 * RH + 0.15, 5.9, 0.38,
    size=9.5, color=SLATE, italic=True)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 9  —  Target Event
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 0.12, 7.5, PRIMARY)

section_label(s, "08  /  COCOK UNTUK", color=TEAL_TEXT)
heading(s, "Platform Ini Tepat\nuntuk Event Anda.", 0.4, 0.62, 9, 1.1, size=34, color=WHITE)

markets = [
    ("Konser & Festival Musik",
     "Event musik skala kecil hingga besar, "
     "festival seni & budaya."),
    ("Seminar & Workshop",
     "Pelatihan profesional, konferensi, "
     "webinar, dan talkshow."),
    ("Event Olahraga",
     "Fun run, tournament, pertandingan "
     "komunitas, dan sport event."),
    ("Pameran & Bazaar",
     "Expo industri, pameran produk, "
     "trade show, dan bazaar UKM."),
    ("Pertunjukan Seni",
     "Teater, pentas tari, stand-up comedy, "
     "dan pameran seni rupa."),
    ("Event Kampus",
     "Acara BEM, UKM, wisuda, ospek, "
     "dan kegiatan mahasiswa."),
]

for i, (title, desc) in enumerate(markets):
    col = i % 3
    row = i // 3
    lx = 0.4 + col * 4.3
    ly = 2.0 + row * 2.55

    rect(s, lx, ly, 4.0, 2.25, DARK2)
    rect(s, lx, ly, 4.0, 0.06,
         PRIMARY if col == 0 else (TEAL_MID if col == 1 else RGBColor(0x2A, 0x6A, 0x6C)))

    num_str = f"0{i+1}"
    txt(s, num_str, lx + 0.2, ly + 0.16, 0.7, 0.4,
        size=9, bold=True, color=PRIMARY, font=F_HEAD)
    txt(s, title, lx + 0.2, ly + 0.62, 3.6, 0.5,
        size=13, bold=True, color=WHITE, font=F_HEAD)
    txt(s, desc, lx + 0.2, ly + 1.2, 3.6, 0.9,
        size=11.5, color=TEAL_TEXT)

foot(s)


# ══════════════════════════════════════════════════════════════
# SLIDE 10  —  Penutup (redesain — tanpa tombol)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[BLANK])

# Full dark background
rect(s, 0, 0, 13.33, 7.5, DARK)

# ── Lime accent bar kiri (tipis, full height) ─────────────────
rect(s, 0, 0, 0.18, 7.5, PRIMARY)

# ── Kiri: statement penutup — center vertikal ─────────────────
# Blok kiri: label(0.28) + gap(0.18) + heading(2.8) + gap(0.3) + subline(0.38)
# Total = 3.94  →  margin = (7.22-3.94)/2 = 1.64
L_X = 0.55
L_W = 6.5

txt(s, "09  /  PENUTUP", L_X, 1.64, L_W, 0.28,
    size=8.5, bold=True, color=TEAL_MID, font=F_HEAD)

txt(s, "Mari Mulai\nBersama.", L_X, 2.1, L_W, 2.8,
    size=62, bold=True, color=WHITE, font=F_HEAD)

rect(s, L_X, 4.95, 1.8, 0.05, PRIMARY)

txt(s, "Platform ticketing online untuk event organizer\nIndonesia — jual tiket mulai hari ini.",
    L_X, 5.12, L_W, 0.72,
    size=13, color=TEAL_TEXT)

# ── Kanan: kontak card — center vertikal ─────────────────────
# Blok kanan: header(0.35) + rule(0.05) + 4 item×1.12 + padding
# Total ≈ 5.28  →  margin = (7.22-5.28)/2 = 0.97
R_X    = 7.65
R_W    = 5.3
R_TOP  = 0.97

# Card background
rect(s, R_X - 0.35, 0, R_W + 0.35, 7.5, DARK2)
rect(s, R_X - 0.35, 0, 0.05, 7.5, TEAL_MID)

txt(s, "HUBUNGI KAMI", R_X, R_TOP, R_W, 0.35,
    size=9, bold=True, color=TEAL_MID, font=F_HEAD)

rect(s, R_X, R_TOP + 0.42, R_W - 0.3, 0.04, TEAL_MID)

contacts = [
    ("WHATSAPP",  "0877-5058-1589"),
    ("EMAIL",     "hello@ingate.id"),
    ("WEBSITE",   "ingate.id"),
    ("LOKASI",    "Ponorogo, Jawa Timur"),
]
for i, (label, value) in enumerate(contacts):
    iy = R_TOP + 0.6 + i * 1.12
    # Lime accent dot
    rect(s, R_X, iy + 0.3, 0.08, 0.08, PRIMARY)
    txt(s, label, R_X + 0.2, iy + 0.22, R_W, 0.3,
        size=8, bold=True, color=TEAL_MID, font=F_HEAD)
    txt(s, value, R_X + 0.2, iy + 0.55, R_W - 0.2, 0.42,
        size=14, bold=True, color=WHITE, font=F_HEAD)
    if i < len(contacts) - 1:
        rect(s, R_X, iy + 1.04, R_W - 0.3, 0.02,
             RGBColor(0x1E, 0x45, 0x48))

foot(s)


# ── Save ──────────────────────────────────────────────────────
out = "/Users/macbookpro/Herd/event/INGATE_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved  →  {out}")
print(f"Slides :  {len(prs.slides)}")
